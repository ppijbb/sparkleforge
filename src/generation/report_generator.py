#!/usr/bin/env python3
"""
Advanced Report Generator for Local Researcher

This module provides comprehensive report generation capabilities including
multiple formats, templates, and automated content synthesis.
"""

import asyncio
import logging
from typing import Dict, Any, List, Optional, Tuple
from datetime import datetime
import json
import os
from pathlib import Path
import jinja2
from jinja2 import Environment, FileSystemLoader
import markdown
import pdfkit
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH

# PowerPoint imports
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    Presentation = None

from src.utils.config_manager import ConfigManager
from src.utils.logger import setup_logger

logger = setup_logger("report_generator", log_level="INFO")


class ReportGenerator:
    """Advanced report generator with multiple output formats."""
    
    def __init__(self, config_path: Optional[str] = None):
        """Initialize the report generator.
        
        Args:
            config_path: Path to configuration file
        """
        self.config_path = config_path
        self.config_manager = ConfigManager(config_path)
        
        # Report settings
        self.templates_dir = Path(self.config_manager.get('templates.directory', './templates'))
        self.output_dir = Path(self.config_manager.get('output.directory', './outputs'))
        self.output_dir.mkdir(exist_ok=True)
        
        # Initialize Jinja2 environment
        self.jinja_env = Environment(
            loader=FileSystemLoader(str(self.templates_dir)),
            autoescape=True
        )
        
        logger.info("Report Generator initialized")
    
    async def generate_research_report(self, 
                                     research_data: Dict[str, Any],
                                     report_type: str = "comprehensive",
                                     output_format: str = "pdf") -> str:
        """Generate a comprehensive research report.
        
        Args:
            research_data: Research data and results
            report_type: Type of report (executive, detailed, academic, presentation)
            output_format: Output format (pdf, html, docx, pptx, markdown)
            
        Returns:
            Path to generated report file
        """
        try:
            logger.info(f"Generating {report_type} report in {output_format} format")
            
            # Prepare report data
            report_data = await self._prepare_report_data(research_data, report_type)
            
            # Generate content based on type
            if report_type == "executive":
                content = await self._generate_executive_summary(report_data)
            elif report_type == "detailed":
                content = await self._generate_detailed_report(report_data)
            elif report_type == "academic":
                content = await self._generate_academic_paper(report_data)
            elif report_type == "presentation":
                content = await self._generate_presentation_slides(report_data)
            else:
                content = await self._generate_comprehensive_report(report_data)
            
            # Generate file based on format
            if output_format == "pdf":
                file_path = await self._generate_pdf_report(content, report_data)
            elif output_format == "html":
                file_path = await self._generate_html_report(content, report_data)
            elif output_format == "docx":
                file_path = await self._generate_docx_report(content, report_data)
            elif output_format == "pptx":
                file_path = await self._generate_pptx_report(content, report_data)
            elif output_format == "markdown":
                file_path = await self._generate_markdown_report(content, report_data)
            else:
                raise ValueError(f"Unsupported output format: {output_format}")
            
            logger.info(f"Report generated successfully: {file_path}")
            return file_path
            
        except Exception as e:
            logger.error(f"Failed to generate research report: {e}")
            raise
    
    async def _prepare_report_data(self, research_data: Dict[str, Any], report_type: str) -> Dict[str, Any]:
        """Prepare data for report generation."""
        try:
            # Extract key information
            objectives = research_data.get('analyzed_objectives', [])
            tasks = research_data.get('decomposed_tasks', [])
            results = research_data.get('execution_results', [])
            evaluation = research_data.get('evaluation_results', {})
            validation = research_data.get('validation_results', {})
            synthesis = research_data.get('final_synthesis', {})
            
            # Calculate metrics
            total_tasks = len(tasks)
            completed_tasks = len([t for t in tasks if t.get('status') == 'completed'])
            success_rate = (completed_tasks / total_tasks * 100) if total_tasks > 0 else 0
            
            quality_score = evaluation.get('overall_score', 0)
            validation_score = validation.get('validation_score', 0)
            
            # Prepare report data
            report_data = {
                'metadata': {
                    'title': research_data.get('user_request', 'Research Report'),
                    'generated_at': datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
                    'report_type': report_type,
                    'objective_id': research_data.get('objective_id', 'unknown')
                },
                'summary': {
                    'total_objectives': len(objectives),
                    'total_tasks': total_tasks,
                    'completed_tasks': completed_tasks,
                    'success_rate': success_rate,
                    'quality_score': quality_score,
                    'validation_score': validation_score
                },
                'objectives': objectives,
                'tasks': tasks,
                'results': results,
                'evaluation': evaluation,
                'validation': validation,
                'synthesis': synthesis,
                'recommendations': self._extract_recommendations(evaluation, validation),
                'key_findings': self._extract_key_findings(results, synthesis)
            }
            
            return report_data
            
        except Exception as e:
            logger.error(f"Failed to prepare report data: {e}")
            raise
    
    async def _generate_executive_summary(self, report_data: Dict[str, Any]) -> str:
        """Generate executive summary content."""
        template = self.jinja_env.get_template('executive_summary.j2')
        return template.render(report_data)
    
    async def _generate_detailed_report(self, report_data: Dict[str, Any]) -> str:
        """Generate detailed report content."""
        template = self.jinja_env.get_template('detailed_report.j2')
        return template.render(report_data)
    
    async def _generate_academic_paper(self, report_data: Dict[str, Any]) -> str:
        """Generate academic paper content."""
        template = self.jinja_env.get_template('academic_paper.j2')
        return template.render(report_data)
    
    async def _generate_presentation_slides(self, report_data: Dict[str, Any]) -> str:
        """Generate presentation slides content."""
        template = self.jinja_env.get_template('presentation_slides.j2')
        return template.render(report_data)
    
    async def _generate_comprehensive_report(self, report_data: Dict[str, Any]) -> str:
        """Generate comprehensive report content."""
        template = self.jinja_env.get_template('comprehensive_report.j2')
        return template.render(report_data)
    
    async def _generate_pdf_report(self, content: str, report_data: Dict[str, Any]) -> str:
        """Generate PDF report."""
        try:
            # Convert markdown to HTML
            html_content = markdown.markdown(content, extensions=['tables', 'codehilite'])
            
            # Add CSS styling
            styled_html = f"""
            <!DOCTYPE html>
            <html>
            <head>
                <meta charset="utf-8">
                <title>{report_data['metadata']['title']}</title>
                <style>
                    body {{ font-family: Arial, sans-serif; margin: 40px; }}
                    h1 {{ color: #2c3e50; border-bottom: 2px solid #3498db; }}
                    h2 {{ color: #34495e; }}
                    h3 {{ color: #7f8c8d; }}
                    table {{ border-collapse: collapse; width: 100%; }}
                    th, td {{ border: 1px solid #ddd; padding: 8px; text-align: left; }}
                    th {{ background-color: #f2f2f2; }}
                    .summary {{ background-color: #f8f9fa; padding: 20px; border-radius: 5px; }}
                    .recommendations {{ background-color: #e8f5e8; padding: 20px; border-radius: 5px; }}
                </style>
            </head>
            <body>
                {html_content}
            </body>
            </html>
            """
            
            # Generate filename
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            filename = f"research_report_{timestamp}.pdf"
            file_path = self.output_dir / filename
            
            # Convert HTML to PDF
            options = {
                'page-size': 'A4',
                'margin-top': '0.75in',
                'margin-right': '0.75in',
                'margin-bottom': '0.75in',
                'margin-left': '0.75in',
                'encoding': "UTF-8",
                'no-outline': None
            }
            
            pdfkit.from_string(styled_html, str(file_path), options=options)
            return str(file_path)
            
        except Exception as e:
            logger.error(f"Failed to generate PDF report: {e}")
            raise
    
    async def _generate_html_report(self, content: str, report_data: Dict[str, Any]) -> str:
        """Generate HTML report."""
        try:
            # Convert markdown to HTML
            html_content = markdown.markdown(content, extensions=['tables', 'codehilite'])
            
            # Add CSS styling
            styled_html = f"""
            <!DOCTYPE html>
            <html>
            <head>
                <meta charset="utf-8">
                <title>{report_data['metadata']['title']}</title>
                <style>
                    body {{ font-family: Arial, sans-serif; margin: 40px; line-height: 1.6; }}
                    h1 {{ color: #2c3e50; border-bottom: 2px solid #3498db; }}
                    h2 {{ color: #34495e; }}
                    h3 {{ color: #7f8c8d; }}
                    table {{ border-collapse: collapse; width: 100%; margin: 20px 0; }}
                    th, td {{ border: 1px solid #ddd; padding: 12px; text-align: left; }}
                    th {{ background-color: #f2f2f2; font-weight: bold; }}
                    .summary {{ background-color: #f8f9fa; padding: 20px; border-radius: 5px; margin: 20px 0; }}
                    .recommendations {{ background-color: #e8f5e8; padding: 20px; border-radius: 5px; margin: 20px 0; }}
                    .key-findings {{ background-color: #fff3cd; padding: 20px; border-radius: 5px; margin: 20px 0; }}
                    .metadata {{ background-color: #e9ecef; padding: 15px; border-radius: 5px; font-size: 0.9em; }}
                </style>
            </head>
            <body>
                {html_content}
            </body>
            </html>
            """
            
            # Generate filename
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            filename = f"research_report_{timestamp}.html"
            file_path = self.output_dir / filename
            
            with open(file_path, 'w', encoding='utf-8') as f:
                f.write(styled_html)
            
            return str(file_path)
            
        except Exception as e:
            logger.error(f"Failed to generate HTML report: {e}")
            raise
    
    async def _generate_docx_report(self, content: str, report_data: Dict[str, Any]) -> str:
        """Generate Word document report."""
        try:
            # Create new document
            doc = Document()
            
            # Add title
            title = doc.add_heading(report_data['metadata']['title'], 0)
            title.alignment = WD_ALIGN_PARAGRAPH.CENTER
            
            # Add metadata
            doc.add_paragraph(f"Generated: {report_data['metadata']['generated_at']}")
            doc.add_paragraph(f"Report Type: {report_data['metadata']['report_type']}")
            doc.add_paragraph(f"Objective ID: {report_data['metadata']['objective_id']}")
            
            # Add summary section
            doc.add_heading('Executive Summary', level=1)
            summary = report_data['summary']
            doc.add_paragraph(f"Total Objectives: {summary['total_objectives']}")
            doc.add_paragraph(f"Total Tasks: {summary['total_tasks']}")
            doc.add_paragraph(f"Completed Tasks: {summary['completed_tasks']}")
            doc.add_paragraph(f"Success Rate: {summary['success_rate']:.1f}%")
            doc.add_paragraph(f"Quality Score: {summary['quality_score']:.2f}")
            doc.add_paragraph(f"Validation Score: {summary['validation_score']:.2f}")
            
            # Add objectives section
            if report_data['objectives']:
                doc.add_heading('Research Objectives', level=1)
                for i, objective in enumerate(report_data['objectives'], 1):
                    doc.add_heading(f"Objective {i}", level=2)
                    doc.add_paragraph(f"Description: {objective.get('description', 'N/A')}")
                    doc.add_paragraph(f"Priority: {objective.get('priority', 'N/A')}")
                    doc.add_paragraph(f"Status: {objective.get('status', 'N/A')}")
            
            # Add results section
            if report_data['results']:
                doc.add_heading('Research Results', level=1)
                for i, result in enumerate(report_data['results'], 1):
                    doc.add_heading(f"Result {i}", level=2)
                    doc.add_paragraph(f"Agent: {result.get('agent', 'N/A')}")
                    doc.add_paragraph(f"Status: {result.get('status', 'N/A')}")
                    if result.get('summary'):
                        doc.add_paragraph(f"Summary: {result['summary']}")
            
            # Add recommendations
            if report_data['recommendations']:
                doc.add_heading('Recommendations', level=1)
                for i, rec in enumerate(report_data['recommendations'], 1):
                    doc.add_paragraph(f"{i}. {rec}")
            
            # Add key findings
            if report_data['key_findings']:
                doc.add_heading('Key Findings', level=1)
                for i, finding in enumerate(report_data['key_findings'], 1):
                    doc.add_paragraph(f"{i}. {finding}")
            
            # Generate filename
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            filename = f"research_report_{timestamp}.docx"
            file_path = self.output_dir / filename
            
            doc.save(str(file_path))
            return str(file_path)
            
        except Exception as e:
            logger.error(f"Failed to generate DOCX report: {e}")
            raise
    
    async def _generate_pptx_report(self, content: str, report_data: Dict[str, Any]) -> str:
        """Generate PowerPoint presentation report."""
        try:
            if not PPTX_AVAILABLE:
                raise ImportError("python-pptx is required for PPTX generation. Install it with: pip install python-pptx")
            
            # Create presentation
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            # Title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = report_data['metadata']['title']
            subtitle.text = f"Generated: {report_data['metadata']['generated_at']}"
            
            # Summary slide
            summary_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(summary_slide_layout)
            title = slide.shapes.title
            title.text = "Executive Summary"
            
            content_shape = slide.placeholders[1]
            tf = content_shape.text_frame
            tf.text = f"Total Objectives: {report_data['summary']['total_objectives']}"
            
            p = tf.add_paragraph()
            p.text = f"Total Tasks: {report_data['summary']['total_tasks']}"
            
            p = tf.add_paragraph()
            p.text = f"Completed Tasks: {report_data['summary']['completed_tasks']}"
            
            p = tf.add_paragraph()
            p.text = f"Success Rate: {report_data['summary']['success_rate']:.1f}%"
            
            p = tf.add_paragraph()
            p.text = f"Quality Score: {report_data['summary']['quality_score']:.2f}"
            
            # Objectives slides
            if report_data['objectives']:
                for i, objective in enumerate(report_data['objectives'], 1):
                    slide = prs.slides.add_slide(summary_slide_layout)
                    title = slide.shapes.title
                    title.text = f"Objective {i}"
                    
                    content_shape = slide.placeholders[1]
                    tf = content_shape.text_frame
                    tf.text = f"Description: {objective.get('description', 'N/A')}"
                    
                    p = tf.add_paragraph()
                    p.text = f"Priority: {objective.get('priority', 'N/A')}"
                    
                    p = tf.add_paragraph()
                    p.text = f"Status: {objective.get('status', 'N/A')}"
            
            # Results slides
            if report_data['results']:
                for i, result in enumerate(report_data['results'], 1):
                    slide = prs.slides.add_slide(summary_slide_layout)
                    title = slide.shapes.title
                    title.text = f"Result {i}"
                    
                    content_shape = slide.placeholders[1]
                    tf = content_shape.text_frame
                    tf.text = f"Agent: {result.get('agent', 'N/A')}"
                    
                    p = tf.add_paragraph()
                    p.text = f"Status: {result.get('status', 'N/A')}"
                    
                    if result.get('summary'):
                        p = tf.add_paragraph()
                        p.text = f"Summary: {result['summary']}"
            
            # Recommendations slide
            if report_data['recommendations']:
                slide = prs.slides.add_slide(summary_slide_layout)
                title = slide.shapes.title
                title.text = "Recommendations"
                
                content_shape = slide.placeholders[1]
                tf = content_shape.text_frame
                for i, rec in enumerate(report_data['recommendations'], 1):
                    if i == 1:
                        tf.text = f"{i}. {rec}"
                    else:
                        p = tf.add_paragraph()
                        p.text = f"{i}. {rec}"
            
            # Key Findings slide
            if report_data['key_findings']:
                slide = prs.slides.add_slide(summary_slide_layout)
                title = slide.shapes.title
                title.text = "Key Findings"
                
                content_shape = slide.placeholders[1]
                tf = content_shape.text_frame
                for i, finding in enumerate(report_data['key_findings'], 1):
                    if i == 1:
                        tf.text = f"{i}. {finding}"
                    else:
                        p = tf.add_paragraph()
                        p.text = f"{i}. {finding}"
            
            # Generate filename
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            filename = f"research_report_{timestamp}.pptx"
            file_path = self.output_dir / filename
            
            prs.save(str(file_path))
            return str(file_path)
            
        except Exception as e:
            logger.error(f"Failed to generate PPTX report: {e}")
            raise
    
    async def _generate_markdown_report(self, content: str, report_data: Dict[str, Any]) -> str:
        """Generate Markdown report."""
        try:
            # Generate filename
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            filename = f"research_report_{timestamp}.md"
            file_path = self.output_dir / filename
            
            with open(file_path, 'w', encoding='utf-8') as f:
                f.write(content)
            
            return str(file_path)
            
        except Exception as e:
            logger.error(f"Failed to generate Markdown report: {e}")
            raise
    
    def _extract_recommendations(self, evaluation: Dict[str, Any], validation: Dict[str, Any]) -> List[str]:
        """Extract recommendations from evaluation and validation data."""
        recommendations = []
        
        # From evaluation
        if evaluation.get('recommendations'):
            recommendations.extend(evaluation['recommendations'])
        
        # From validation
        if validation.get('recommendations'):
            recommendations.extend(validation['recommendations'])
        
        # Generate based on scores
        quality_score = evaluation.get('overall_score', 0)
        if quality_score < 0.7:
            recommendations.append("Improve research quality through more thorough data collection")
        
        validation_score = validation.get('validation_score', 0)
        if validation_score < 0.8:
            recommendations.append("Enhance result validation process")
        
        return recommendations
    
    def _extract_key_findings(self, results: List[Dict[str, Any]], synthesis: Dict[str, Any]) -> List[str]:
        """Extract key findings from results and synthesis."""
        findings = []
        
        # From synthesis
        if synthesis.get('key_findings'):
            findings.extend(synthesis['key_findings'])
        
        # From results
        for result in results:
            if result.get('key_findings'):
                findings.extend(result['key_findings'])
        
        return findings
